"""Assemble EVERY project figure into a single PowerPoint deck, lightly organised.

Globs all `HCV1_*.png` in the vault attachments plus the spatial-sweep PNGs in
`cca/figures/`, groups them into logical sections (one section-divider slide each),
and lays one figure per slide scaled to fit (aspect preserved), with a prettified
title and the raw filename as a caption for traceability. Nothing is dropped — any
figure not matched by a section falls into "Other (unsorted)" at the end.

Reproducible: re-run after the bin50 / FS-incl robustness figures land and the deck
picks them up automatically. Dependency-light — python-pptx + stdlib only (no Pillow;
PNG dimensions are read from the file header). Run:
    uv run --with python-pptx python cca/scripts/make_figure_deck.py
"""

from __future__ import annotations

import struct
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

ATT = Path.home() / "Documents" / "ResearchVault" / "attachments"
SWEEPS = Path(__file__).resolve().parents[1] / "figures"
OUT = (Path.home() / "Documents" / "ResearchVault" / "Projects" / "Hippocampus-V1"
       / "HCV1_all_figures.pptx")

EMU_PER_PX = 9525  # 96 dpi
INK = RGBColor(0x23, 0x23, 0x1F)
MUTED = RGBColor(0x7A, 0x79, 0x72)
ACCENT = RGBColor(0x1F, 0x9E, 0x8F)

# Figures to drop entirely (sweeps removed by request; KCCA kept to the verdict only).
def SKIP(s):
    if "sweep" in s:
        return True
    if "kcca" in s and "vs_linear" not in s and "early" not in s:
        return True
    return False


# Ordered sections: (title, subtitle, predicate on lowercase filename stem).
# First match wins, so put the specific (KCCA, early) before the generic.
SECTIONS = [
    ("Graphical abstract", "the one-slide story",
     lambda s: "graphical_abstract" in s),
    ("Statistics at a glance", "summary tables — learners vs pooled, colour-coded",
     lambda s: "stats_" in s),
    ("Levels — the coupling hierarchy", "held-out CC / n_sig / IFI per area pair",
     lambda s: "_levels" in s and "kcca" not in s),
    ("Strength & sparsity vs learning", "trajectory, animals-as-n (25 ms unless noted)",
     lambda s: any(k in s for k in ("gini_traj", "slopes_", "gini_control",
                                    "learnervsnon")) and "kcca" not in s),
    ("Directionality (IFI)", "trajectory, per-epoch, and the held-out window sweep",
     lambda s: any(k in s for k in ("direction", "ifi_traj", "ifi_windows"))
     and "kcca" not in s and "early" not in s),
    ("Dims-as-n — the Buzsáki unit", "animals-as-n vs significant-dimensions-as-n",
     lambda s: any(k in s for k in ("units_", "trajdims", "_epochs")) and "kcca" not in s),
    ("Subspace rotation / reorientation", "cross-window angle vs split-half noise floor",
     lambda s: "rotation" in s and "kcca" not in s),
    ("Transition — uncued → cued", "within-session block contrast",
     lambda s: "transition" in s and "kcca" not in s),
    ("Nonlinear (kernel) CCA — verdict", "largely linear (KCCA ≈ linear); rest dropped for now",
     lambda s: "kcca" in s and "vs_linear" in s and "early" not in s),
    ("The first ~10 trials", "very-early-trial readouts (block & projected)",
     lambda s: "early" in s),
]

PRETTY = {
    "fsexcl": "FS-excl", "fsincl": "FS-incl", "bin25": "25 ms", "bin50": "50 ms",
    "ifi": "IFI", "cc1": "CC1", "cc": "CC", "kcca": "KCCA", "cca": "CCA",
    "gini": "Gini", "traj": "trajectory", "proj": "projected", "rsc": "RSC",
    "v1": "V1", "ca1": "CA1", "lp": "LP", "vs": "vs", "n": "n",
}


def png_size(path: Path):
    with open(path, "rb") as f:
        head = f.read(24)
    if head[:8] != b"\x89PNG\r\n\x1a\n":
        return 1600, 1000
    w, h = struct.unpack(">II", head[16:24])
    return w, h


def prettify(stem: str) -> str:
    body = stem.replace("HCV1_", "")
    parts = [PRETTY.get(p, p) for p in body.split("_") if p]
    return " ".join(parts)


def collect():
    figs = [f for f in sorted(ATT.glob("HCV1_*.png")) if not SKIP(f.stem.lower())]
    assigned, sectioned = set(), [(t, sub, []) for t, sub, _ in SECTIONS]
    for f in figs:
        s = f.stem.lower()
        for i, (_t, _sub, pred) in enumerate(SECTIONS):
            if pred(s):
                sectioned[i][2].append(f)
                assigned.add(f)
                break
    leftover = [f for f in figs if f not in assigned]
    if leftover:
        sectioned.append(("Other (unsorted)", "anything not matched above", leftover))
    return [(t, sub, fs) for t, sub, fs in sectioned if fs]


def _textbox(slide, left, top, width, height, text, size, color, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = size
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]
    sections = collect()
    total = sum(len(fs) for _, _, fs in sections)

    # Title slide
    s = prs.slides.add_slide(blank)
    _textbox(s, Inches(0.8), Inches(2.6), SW - Inches(1.6), Inches(1.5),
             "Hippocampal–cortical communication subspaces", Pt(34), INK, bold=True,
             align=PP_ALIGN.CENTER)
    _textbox(s, Inches(0.8), Inches(3.9), SW - Inches(1.6), Inches(1.2),
             f"All figures — {total} panels across {len(sections)} sections "
             "(both FS conditions, animals- & dims-as-n, linear & kernel CCA)",
             Pt(16), MUTED, align=PP_ALIGN.CENTER)

    for sec_i, (title, sub, figs) in enumerate(sections, 1):
        # divider
        d = prs.slides.add_slide(blank)
        bar = d.shapes.add_shape(1, Inches(0.8), Inches(3.0), Inches(0.18), Inches(1.4))
        bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
        _textbox(d, Inches(1.2), Inches(3.0), SW - Inches(2.0), Inches(1.0),
                 f"{sec_i}.  {title}", Pt(30), INK, bold=True)
        _textbox(d, Inches(1.2), Inches(4.0), SW - Inches(2.0), Inches(0.8),
                 f"{sub}   ·   {len(figs)} figure{'s' if len(figs) != 1 else ''}",
                 Pt(15), MUTED)
        # one figure per slide
        for f in figs:
            sl = prs.slides.add_slide(blank)
            _textbox(sl, Inches(0.4), Inches(0.18), SW - Inches(0.8), Inches(0.6),
                     f"{title}  —  {prettify(f.stem)}", Pt(17), INK, bold=True)
            iw, ih = png_size(f)
            maxw, maxh = SW - Inches(0.8), SH - Inches(1.4)
            scale = min(maxw / (iw * EMU_PER_PX), maxh / (ih * EMU_PER_PX))
            w = Emu(int(iw * EMU_PER_PX * scale))
            h = Emu(int(ih * EMU_PER_PX * scale))
            left = Emu(int((SW - w) / 2))
            top = Emu(int(Inches(0.95) + (maxh - h) / 2))
            sl.shapes.add_picture(str(f), left, top, width=w, height=h)
            _textbox(sl, Inches(0.4), SH - Inches(0.34), SW - Inches(0.8), Inches(0.3),
                     f.name, Pt(9), MUTED, align=PP_ALIGN.CENTER)

    prs.save(str(OUT))
    print(f"wrote {OUT}")
    print(f"  {len(prs.slides._sldIdLst)} slides, {total} figures, {len(sections)} sections:")
    for t, _sub, fs in sections:
        print(f"    {len(fs):>2}  {t}: {', '.join(f.stem.replace('HCV1_', '') for f in fs)}")


if __name__ == "__main__":
    main()
