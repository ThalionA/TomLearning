"""Assemble ONLY the smoothed-10 ms figures into a single PowerPoint deck.

Every smoothed-10 ms figure carries `bin10` in its filename, so a `HCV1_*bin10*.png`
glob selects exactly the current smoothed set (both FS, the new lag-CC / paired /
per-area-Gini / weight-CDF / 5-trial-bin figures included) and nothing from the 25 ms
or KCCA suites. One figure per slide, scaled to fit, prettified title + filename caption;
result-themed section dividers. Re-run to pick up figures as later runs land.

Dependency-light (python-pptx + stdlib; PNG dims from the header, no Pillow). Run:
    uv run --with python-pptx python cca/scripts/make_smoothed_deck.py
"""
from __future__ import annotations

import struct
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

ATT = Path.home() / "Documents" / "ResearchVault" / "attachments"
OUT = (Path.home() / "Documents" / "ResearchVault" / "Projects" / "Hippocampus-V1"
       / "HCV1_smoothed_10ms_results.pptx")

EMU_PER_PX = 9525  # 96 dpi
INK = RGBColor(0x23, 0x23, 0x1F)
MUTED = RGBColor(0x7A, 0x79, 0x72)
ACCENT = RGBColor(0x1F, 0x9E, 0x8F)

# Ordered sections: (title, subtitle, predicate on lowercase stem). First match wins,
# so area-Gini/CDF precede the generic "transition" bucket, paired-epochs precede it too.
SECTIONS = [
    ("Coupling levels — the hierarchy", "held-out CC1 / n_sig / IFI per area pair",
     lambda s: "_levels" in s),
    ("Directionality & lagged-CC curves",
     "IFI over learning + held-out window sweep + the lag-CC profiles (new) + stats",
     lambda s: any(k in s for k in ("direction", "ifi_traj", "ifi_windows", "lagcc",
                                    "stats_directionality"))),
    ("Strength & sparsity vs learning",
     "CC1 / Gini trajectories, slopes, learner-vs-non, CCA-independent control",
     lambda s: any(k in s for k in ("gini_traj", "slopes_", "gini_control",
                                    "learnervsnon", "stats_slopes"))),
    ("Per-dimension trajectory", "significant-dimension CCs across learning",
     lambda s: "trajdims" in s),
    ("Subspace rotation / reorientation", "cross-window angle vs split-half noise floor",
     lambda s: "rotation" in s),
    ("Paired: naïve vs trained (learners)",
     "per-animal paired contrasts — CC1 / n_sig / IFI / Gini / mincc, FS-incl",
     lambda s: "paired_epochs" in s),
    ("Per-area Gini & weight distribution",
     "Gini per area across all partners + Fig-4b Lorenz weight CDF",
     lambda s: "areagini" in s or "weightcdf" in s),
    ("Transition — uncued → cued",
     "all-relationship map + per-animal paired contrasts",
     lambda s: "transition" in s),
    ("5-trial-bin trajectories (early learning)",
     "metrics across trial bins 1-5…26-30, learners vs non-learners",
     lambda s: "trajbins" in s),
]

PRETTY = {
    "fsexcl": "FS-excl", "fsincl": "FS-incl", "bin10": "10 ms", "pooled": "(all animals)",
    "ifi": "IFI", "cc1": "CC1", "cc": "CC", "cca": "CCA", "gini": "Gini",
    "traj": "trajectory", "trajdims": "per-dim trajectory", "trajbins": "5-trial bins",
    "lagcc": "lag-CC curves", "areagini": "per-area Gini", "weightcdf": "weight CDF",
    "paired": "paired", "epochs": "naïve-vs-trained", "transition": "uncued-vs-cued",
    "learnervsnon": "learner-vs-non", "learn": "learners", "all": "all animals",
    "mincc": "mincc", "rsc": "RSC", "v1": "V1", "ca1": "CA1", "lp": "LP", "vs": "vs",
    "n": "n", "sig": "sig", "rotation": "rotation", "direction": "direction",
}


def png_size(path: Path):
    with open(path, "rb") as f:
        head = f.read(24)
    if head[:8] != b"\x89PNG\r\n\x1a\n":
        return 1600, 1000
    return struct.unpack(">II", head[16:24])


def prettify(stem: str) -> str:
    body = stem.replace("HCV1_", "")
    return " ".join(PRETTY.get(p, p) for p in body.split("_") if p)


def collect():
    figs = sorted(ATT.glob("HCV1_*bin10*.png"))
    sectioned = [(t, sub, []) for t, sub, _ in SECTIONS]
    assigned = set()
    for f in figs:
        s = f.stem.lower()
        for i, (_t, _sub, pred) in enumerate(SECTIONS):
            if pred(s):
                sectioned[i][2].append(f); assigned.add(f); break
    leftover = [f for f in figs if f not in assigned]
    if leftover:
        sectioned.append(("Other (unsorted)", "not matched above", leftover))
    return [(t, sub, fs) for t, sub, fs in sectioned if fs]


def _textbox(slide, left, top, width, height, text, size, color, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = size; r.font.bold = bold; r.font.color.rgb = color
    return tb


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]
    sections = collect()
    total = sum(len(fs) for _, _, fs in sections)

    s = prs.slides.add_slide(blank)
    _textbox(s, Inches(0.8), Inches(2.5), SW - Inches(1.6), Inches(1.5),
             "Hippocampal–cortical communication subspaces — smoothed 10 ms results",
             Pt(30), INK, bold=True, align=PP_ALIGN.CENTER)
    _textbox(s, Inches(0.8), Inches(3.9), SW - Inches(1.6), Inches(1.4),
             f"Gaussian σ=2.5 ms spike-train smoothing · 10 ms bins · both FS conditions\n"
             f"{total} figures across {len(sections)} sections   ·   held-out, leak-free CV throughout",
             Pt(15), MUTED, align=PP_ALIGN.CENTER)

    for sec_i, (title, sub, figs) in enumerate(sections, 1):
        d = prs.slides.add_slide(blank)
        bar = d.shapes.add_shape(1, Inches(0.8), Inches(3.0), Inches(0.18), Inches(1.4))
        bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
        _textbox(d, Inches(1.2), Inches(3.0), SW - Inches(2.0), Inches(1.0),
                 f"{sec_i}.  {title}", Pt(30), INK, bold=True)
        _textbox(d, Inches(1.2), Inches(4.0), SW - Inches(2.0), Inches(0.8),
                 f"{sub}   ·   {len(figs)} figure{'s' if len(figs) != 1 else ''}",
                 Pt(15), MUTED)
        for f in figs:
            sl = prs.slides.add_slide(blank)
            _textbox(sl, Inches(0.4), Inches(0.18), SW - Inches(0.8), Inches(0.6),
                     f"{title}  —  {prettify(f.stem)}", Pt(16), INK, bold=True)
            iw, ih = png_size(f)
            maxw, maxh = SW - Inches(0.8), SH - Inches(1.4)
            scale = min(maxw / (iw * EMU_PER_PX), maxh / (ih * EMU_PER_PX))
            w = Emu(int(iw * EMU_PER_PX * scale)); h = Emu(int(ih * EMU_PER_PX * scale))
            sl.shapes.add_picture(str(f), Emu(int((SW - w) / 2)),
                                  Emu(int(Inches(0.95) + (maxh - h) / 2)), width=w, height=h)
            _textbox(sl, Inches(0.4), SH - Inches(0.34), SW - Inches(0.8), Inches(0.3),
                     f.name, Pt(9), MUTED, align=PP_ALIGN.CENTER)

    prs.save(str(OUT))
    print(f"wrote {OUT}")
    print(f"  {len(prs.slides._sldIdLst)} slides, {total} figures, {len(sections)} sections:")
    for t, _sub, fs in sections:
        print(f"    {len(fs):>2}  {t}")


if __name__ == "__main__":
    main()
